import os
from typing import Optional
from google.adk.agents import Agent  
from pptx import Presentation 
from pptx.util import Pt  


def save_ppt(
    title: str,
    slides: list[str],
    theme: Optional[str] = None,
    slide_count: Optional[int] = None,
    bullet_font_pt: int = 18  
):
    

    # Ensure base folder exists
    os.makedirs("source", exist_ok=True)

    # Create subfolder for PPT
    ppt_folder = os.path.join("source", "ppt")
    os.makedirs(ppt_folder, exist_ok=True)

    # Create presentation
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"✨ Auto-generated Presentation ({theme or 'default'})"

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title + Content
    max_slides = slide_count if slide_count else len(slides)

    for i, content in enumerate(slides[:max_slides], start=1):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = f"{title} - Part {i}"

        # Add bullet points properly
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for line in content.split("\n"):
            if line.strip():  # Skip empty lines
                p = text_frame.add_paragraph()
                p.text = line.strip().lstrip("*• ")  # remove extra * or • if present
                p.level = 0  # bullet level

                # Apply reduced font size
                if p.runs:
                    for run in p.runs:
                        run.font.size = Pt(bullet_font_pt)

    # Save PPT
    ppt_path = os.path.join(ppt_folder, f"{title}_presentation.pptx")
    prs.save(ppt_path)
    try:
        os.startfile(ppt_path)  # Auto-open (Windows only)
    except Exception:
        pass

    return f"PPT saved at {ppt_path}"


# Agent definition
ppt_agent = Agent(
    name="ppt_agent",
    model="gemini-2.0-flash",
    instruction="""
You are a PowerPoint Builder Agent.

**Goal:** Create presentations quickly and simply.

**Interaction Rules (Voice-Friendly):**
1. Only ask for what’s absolutely needed:
   - **Topic**
   - **Number of slides**
2. If the user gives one but not the other, politely ask for the missing info in a very short question.  
   - Example: "How many slides?" or "What’s the topic?"
3. Never ask for theme or title unless the user explicitly gives it.
   - Default title = topic
   - Default theme = 'default'

**Content Rules:**
- Generate exactly the requested number of slides.
- Each slide must have **5–7 clear bullet points**.
- Keep content **informative and useful** (examples, definitions, advantages, challenges, etc.).

**Response Style:**
- Keep answers short and natural (1 short sentence).
- Examples:
  - Instead of: "Okay, I will now create a PowerPoint presentation for you."
  - Say: "Making your PPT now."
  - After saving: " PPT ready and opened."

**Process:**
1. Collect topic + slide count.
2. Generate structured slide list.
3. Call the `save_ppt` tool with ( title, slides, theme, slide_count).
4. Respond simply after saving (<= 1 sentence).
""",
    tools=[save_ppt]
)
